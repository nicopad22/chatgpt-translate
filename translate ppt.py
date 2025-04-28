from pptx import Presentation
import openai
import time
import os

# VARIABLES
root = "C:\\Users\\nicol\\Downloads\\CHATGPT-TRANSLATE\\Excel\\BHPF 1\\"
output_language = "es"
# output_language = "en"

# OpenAI API client
# Store there your API key
with open(".\\TOKEN.txt") as token_file:
    client = openai.OpenAI(api_key=[line for line in token_file][0])

# general translation function
def translate(text):
    if text.strip() != "":
        if output_language == "es":
            return translate_to_spanish(text)
        elif output_language == "en":
            return translate_to_english(text)
    else:
        return ""

# Function to translate text
def translate_to_spanish(text):
    result = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
                {"role": "system",
                 "content": "Eres un asistente de traducción. Traduces cualquier texto que recibas al español, sin realizar comentarios adicionales."},
                {"role": "user", "content": f"{text.strip()}"}
            ],
        temperature=0.7,
    )

    translation = result.choices[0].message.content
    print(text + " / " + translation)

    return translation

# Function to translate text
def translate_to_english(text):
    result = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
                {"role": "system",
                 "content": "You are a translation assistant. You translate any text into english, without making any comments even if it is not translatable."},
                {"role": "user", "content": f"{text.strip()}"}
            ],
        temperature=0.7,
    )

    translation = result.choices[0].message.content
    print(text + " / " + translation)

    return translation

# List out all files in folder, filter out != .xlsx
files = os.listdir(root)
print(f"files found: ")
ppt_files = []
for file in files:
    if file[-5:] == ".pptx":
        print(" - " + file)
        ppt_files.append(file)

for file in ppt_files:
    newfile = file[:-5] + " [TRANSLATED].pptx"

    t = os.popen(f"copy \"{root + file}\" \"{root + newfile}\"")

    old_ppt = Presentation(root + file)
    new_ppt = Presentation(root + newfile)

    for slide in new_ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        
                        run.text = translate(text)
    
    new_ppt.save(root + newfile)